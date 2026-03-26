# tools/documents.py — Document generation and extraction tools.
#
# All output written to workspace/ only. Path traversal rejected.
# Libraries imported lazily — missing packages return a clear error.
#
# Tools: draft_pptx, draft_docx, draft_xlsx, generate_qr,
#        render_diagram, render_template, financial_calc, pdf_extract

import re
import subprocess
import tempfile
from pathlib import Path

_WORKSPACE = Path(__file__).parent.parent / "workspace"


def _resolve_output_path(filename: str, expected_suffix: str = None):
    """Validate filename, strip workspace/ prefix, check traversal, return Path."""
    clean = filename.strip()
    if clean.startswith("workspace/") or clean.startswith("workspace\\"):
        clean = clean[len("workspace/"):]
    clean = clean.strip()
    if not clean:
        return "Error: filename must not be empty."
    if ".." in Path(clean).parts:
        return "Rejected: path traversal not allowed."
    if expected_suffix and not clean.lower().endswith(expected_suffix):
        return f"Error: filename must end with '{expected_suffix}'."
    out = _WORKSPACE / clean
    out.parent.mkdir(parents=True, exist_ok=True)
    return out


# ── PowerPoint ────────────────────────────────────────────────────────────────

def draft_pptx(filename: str, slides: list) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return "python-pptx not installed — add 'python-pptx' to requirements.txt."

    out = _resolve_output_path(filename, ".pptx")
    if isinstance(out, str):
        return out

    NAVY = RGBColor(31, 73, 125)
    WHITE = RGBColor(255, 255, 255)
    DARK_GREY = RGBColor(51, 51, 51)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, slide_data in enumerate(slides):
        slide_title = str(slide_data.get("title", f"Slide {idx + 1}"))
        slide_content = str(slide_data.get("content", ""))
        slide_notes = str(slide_data.get("notes", ""))
        slide = prs.slides.add_slide(blank_layout)
        is_title_slide = (idx == 0)

        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = NAVY if is_title_slide else WHITE

        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(1.4))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(40 if is_title_slide else 32)
        run.font.bold = True
        run.font.color.rgb = WHITE if is_title_slide else NAVY

        if slide_content.strip():
            content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(12.0), Inches(5.0))
            ctf = content_box.text_frame
            ctf.word_wrap = True
            first = True
            for line in slide_content.splitlines():
                line = line.strip()
                if not line:
                    continue
                cp = ctf.paragraphs[0] if first else ctf.add_paragraph()
                first = False
                if line[0] in ("-", "•", "*"):
                    cp.level = 1
                    line = line.lstrip("-•* ").strip()
                else:
                    cp.level = 0
                cp.text = line
                run = cp.runs[0]
                run.font.name = "Calibri"
                run.font.size = Pt(20 if is_title_slide else 18)
                run.font.color.rgb = WHITE if is_title_slide else DARK_GREY

        if slide_notes.strip():
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_notes

    prs.save(str(out))
    n = len(slides)
    return f"Created {filename} — {n} slide{'s' if n != 1 else ''}."


# ── Word document ─────────────────────────────────────────────────────────────

def draft_docx(filename: str, content: str, title: str = "") -> str:
    try:
        from docx import Document
        from docx.shared import Inches, Pt
    except ImportError:
        return "python-docx not installed — add 'python-docx' to requirements.txt."

    out = _resolve_output_path(filename, ".docx")
    if isinstance(out, str):
        return out

    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

    if title:
        doc.core_properties.title = title

    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(12)

    def _apply_inline(para, text: str) -> None:
        token_re = re.compile(r'(\*\*(.+?)\*\*|\*(.+?)\*)')
        cursor = 0
        for m in token_re.finditer(text):
            plain = text[cursor:m.start()]
            if plain:
                para.add_run(plain)
            if m.group(0).startswith("**"):
                para.add_run(m.group(2)).bold = True
            else:
                para.add_run(m.group(3)).italic = True
            cursor = m.end()
        remainder = text[cursor:]
        if remainder:
            para.add_run(remainder)

    for raw_line in content.splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            doc.add_paragraph("")
            continue
        matched = False
        for prefix in ("### ", "## ", "# "):
            if line.startswith(prefix):
                doc.add_heading(line[len(prefix):], level={"# ": 1, "## ": 2, "### ": 3}[prefix])
                matched = True
                break
        if matched:
            continue
        if line.startswith("- "):
            para = doc.add_paragraph(style="List Bullet")
            _apply_inline(para, line[2:])
            continue
        para = doc.add_paragraph(style="Normal")
        _apply_inline(para, line)

    doc.save(str(out))
    return f"Created {filename}."


# ── Excel workbook ────────────────────────────────────────────────────────────

def draft_xlsx(filename: str, sheets: list) -> str:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter
    except ImportError:
        return "openpyxl not installed — add 'openpyxl' to requirements.txt."

    out = _resolve_output_path(filename, ".xlsx")
    if isinstance(out, str):
        return out

    HEADER_BG = "1F4980"
    HEADER_FG = "FFFFFF"
    ALT_BG = "EEF3FB"

    wb = Workbook()
    wb.remove(wb.active)

    for sheet_def in sheets:
        sheet_name = str(sheet_def.get("name", "Sheet"))[:31]
        headers = list(sheet_def.get("headers", []))
        rows = list(sheet_def.get("rows", []))
        col_widths = sheet_def.get("col_widths", [])

        ws = wb.create_sheet(title=sheet_name)

        if headers:
            for col_idx, header in enumerate(headers, start=1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = Font(name="Calibri", bold=True, color=HEADER_FG)
                cell.fill = PatternFill("solid", fgColor=HEADER_BG)
                cell.alignment = Alignment(horizontal="center", vertical="center")
            last_col = get_column_letter(len(headers))
            ws.auto_filter.ref = f"A1:{last_col}1"
            ws.freeze_panes = "A2"

        for row_idx, row in enumerate(rows, start=2):
            fill = PatternFill("solid", fgColor=ALT_BG) if row_idx % 2 == 0 else None
            for col_idx, value in enumerate(row, start=1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.font = Font(name="Calibri", size=11)
                cell.alignment = Alignment(vertical="center")
                if fill:
                    cell.fill = fill

        for col_idx in range(1, len(headers) + 1):
            col_letter = get_column_letter(col_idx)
            if col_widths and col_idx - 1 < len(col_widths):
                ws.column_dimensions[col_letter].width = col_widths[col_idx - 1]
            else:
                max_len = len(str(headers[col_idx - 1])) if headers else 8
                for row in rows:
                    if col_idx - 1 < len(row):
                        max_len = max(max_len, len(str(row[col_idx - 1])))
                ws.column_dimensions[col_letter].width = min(max_len + 4, 60)
        ws.row_dimensions[1].height = 20

    wb.save(str(out))
    total_rows = sum(len(s.get("rows", [])) for s in sheets)
    return f"Created {filename} — {len(sheets)} sheet(s), {total_rows} data row(s)."


# ── QR code ───────────────────────────────────────────────────────────────────

def generate_qr(filename: str, data: str, size: int = 10) -> str:
    try:
        import qrcode
    except ImportError:
        return "qrcode not installed — add 'qrcode[pil]' to requirements.txt."
    try:
        from PIL import Image  # noqa
    except ImportError:
        return "Pillow not installed — add 'Pillow' to requirements.txt."

    out = _resolve_output_path(filename, ".png")
    if isinstance(out, str):
        return out

    if not data or not data.strip():
        return "Error: 'data' must not be empty."

    size = max(1, min(int(size), 50))
    qr = qrcode.QRCode(version=None, error_correction=qrcode.constants.ERROR_CORRECT_H, box_size=size, border=4)
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    img.save(str(out))
    return f"QR code saved to {filename}."


# ── Diagram rendering ─────────────────────────────────────────────────────────

def render_diagram(filename: str, diagram_code: str, diagram_type: str = "mermaid") -> str:
    out = _resolve_output_path(filename)
    if isinstance(out, str):
        return out

    diagram_type = diagram_type.lower().strip()

    if diagram_type == "mermaid":
        with tempfile.NamedTemporaryFile(mode="w", suffix=".mmd", delete=False, encoding="utf-8") as tmp:
            tmp.write(diagram_code)
            tmp_path = tmp.name
        try:
            result = subprocess.run(["mmdc", "-i", tmp_path, "-o", str(out)], capture_output=True, text=True, timeout=60)
            if result.returncode != 0:
                return f"mmdc error: {(result.stderr or result.stdout or 'unknown error').strip()}"
            return f"Diagram rendered to {filename}."
        except FileNotFoundError:
            return "mmdc not found — install: npm install -g @mermaid-js/mermaid-cli"
        except subprocess.TimeoutExpired:
            return "mmdc timed out (>60s)."
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    elif diagram_type == "graphviz":
        try:
            import graphviz as gv
        except ImportError:
            return "graphviz not installed — add 'graphviz' to requirements.txt."
        suffix = out.suffix.lstrip(".").lower() or "png"
        if suffix not in {"png", "svg", "pdf"}:
            suffix = "png"
        try:
            src = gv.Source(diagram_code, format=suffix)
            rendered = src.render(filename=str(out.with_suffix("")), cleanup=True, quiet=True)
            rendered_path = Path(rendered)
            if rendered_path != out and rendered_path.exists():
                rendered_path.rename(out)
            return f"Diagram rendered to {filename}."
        except Exception as e:
            return f"graphviz error: {e}"
    else:
        return f"Unknown diagram_type '{diagram_type}' — must be 'mermaid' or 'graphviz'."


# ── Jinja2 template ───────────────────────────────────────────────────────────

def render_template(filename: str, template_str: str, variables: dict) -> str:
    try:
        from jinja2 import Environment, BaseLoader, TemplateError
    except ImportError:
        return "jinja2 not installed — add 'jinja2' to requirements.txt."

    out = _resolve_output_path(filename)
    if isinstance(out, str):
        return out

    if not template_str.strip():
        return "Error: template_str must not be empty."

    try:
        env = Environment(loader=BaseLoader(), autoescape=False, keep_trailing_newline=True)
        rendered = env.from_string(template_str).render(**(variables or {}))
    except TemplateError as e:
        return f"Template error: {e}"

    out.write_text(rendered, encoding="utf-8")
    return f"Template rendered to {filename} ({len(rendered)} characters)."


# ── Financial calculations ────────────────────────────────────────────────────

def financial_calc(calc_type: str, params: dict, filename: str = "") -> str:
    try:
        import pandas as pd
        import numpy as np
    except ImportError:
        missing = []
        try: import pandas  # noqa
        except ImportError: missing.append("pandas")
        try: import numpy  # noqa
        except ImportError: missing.append("numpy")
        return f"Libraries not installed — add {missing} to requirements.txt."

    calc_type = (calc_type or "").strip().lower()
    lines = []

    if calc_type == "projection":
        start = float(params.get("start_revenue", 0))
        growth = float(params.get("growth_rate_pct", 0)) / 100.0
        months = int(params.get("months", 12))
        costs = float(params.get("costs_monthly", 0))
        records = []
        cumulative = 0.0
        for m in range(1, months + 1):
            rev = start * ((1 + growth) ** (m - 1))
            profit = rev - costs
            cumulative += profit
            records.append({"Month": m, "Revenue ($)": round(rev, 2), "Costs ($)": round(costs, 2), "Profit ($)": round(profit, 2), "Cumulative Profit": round(cumulative, 2)})
        df = pd.DataFrame(records)
        lines.append("=== Revenue Projection ===")
        lines.append(f"Start: ${start:,.2f}/mo | Growth: {growth*100:.1f}%/mo | Costs: ${costs:,.2f}/mo\n")
        lines.append(df.to_string(index=False))
        lines.append(f"\nFinal revenue: ${records[-1]['Revenue ($)']:,.2f} | Total profit: ${cumulative:,.2f}")

    elif calc_type == "break_even":
        fixed = float(params.get("fixed_costs", 0))
        var = float(params.get("variable_cost_per_unit", 0))
        price = float(params.get("price_per_unit", 0))
        monthly = float(params.get("monthly_sales", 0))
        if price <= var:
            return "Error: price_per_unit must be greater than variable_cost_per_unit."
        margin = price - var
        be_units = fixed / margin
        be_revenue = be_units * price
        lines.append("=== Break-Even Analysis ===")
        lines.append(f"Fixed costs:          ${fixed:,.2f}")
        lines.append(f"Variable cost/unit:   ${var:,.2f}")
        lines.append(f"Price/unit:           ${price:,.2f}")
        lines.append(f"Contribution margin:  ${margin:,.2f}/unit")
        lines.append(f"Break-even units:     {be_units:,.1f}")
        lines.append(f"Break-even revenue:   ${be_revenue:,.2f}")
        if monthly > 0:
            lines.append(f"Months to break even: {be_units / monthly:.1f}")

    elif calc_type == "market_size":
        tam_units = int(params.get("total_addressable", 0))
        capture_pct = float(params.get("capture_rate_pct", 0)) / 100.0
        price = float(params.get("price_per_customer", 0))
        sam_units = tam_units * capture_pct
        lines.append("=== Market Size Estimate ===")
        lines.append(f"TAM: {tam_units:,} customers / ${tam_units * price:,.0f}")
        lines.append(f"Capture rate: {capture_pct * 100:.1f}%")
        lines.append(f"SAM: {sam_units:,.0f} customers / ${sam_units * price:,.0f}")
        lines.append(f"ARR estimate: ${sam_units * price:,.0f}")

    elif calc_type == "roi":
        investment = float(params.get("investment", 0))
        monthly_return = float(params.get("monthly_return", 0))
        months = int(params.get("months", 12))
        if investment <= 0:
            return "Error: investment must be > 0."
        total_return = monthly_return * months
        roi_pct = ((total_return - investment) / investment) * 100.0
        payback = investment / monthly_return if monthly_return > 0 else float("inf")
        monthly_rate = (1 + 0.10) ** (1 / 12) - 1
        cash_flows = np.array([monthly_return] * months)
        periods = np.arange(1, months + 1)
        npv = float(np.sum(cash_flows / (1 + monthly_rate) ** periods) - investment)
        lines.append("=== ROI Analysis ===")
        lines.append(f"Investment:       ${investment:,.2f}")
        lines.append(f"Monthly return:   ${monthly_return:,.2f}")
        lines.append(f"Duration:         {months} months")
        lines.append(f"Total return:     ${total_return:,.2f}")
        lines.append(f"ROI:              {roi_pct:.1f}%")
        lines.append(f"Payback:          {'never' if payback == float('inf') else f'{payback:.1f} months'}")
        lines.append(f"NPV (10% annual): ${npv:,.2f}")

    else:
        return f"Unknown calc_type '{calc_type}'. Valid: projection, break_even, market_size, roi."

    output = "\n".join(lines)

    if filename and filename.strip():
        out = _resolve_output_path(filename)
        if not isinstance(out, str):
            out.write_text(output, encoding="utf-8")
            return output + f"\n\nSaved to {filename}."

    return output


# ── PDF extract ───────────────────────────────────────────────────────────────

def pdf_extract(path: str) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        return "pypdf not installed — add 'pypdf' to requirements.txt."

    clean = path.strip()
    if clean.startswith("workspace/"):
        clean = clean[len("workspace/"):]
    if ".." in Path(clean).parts:
        return "Rejected: path traversal not allowed."

    file_path = _WORKSPACE / clean
    if not file_path.exists():
        return f"File not found: workspace/{clean}"
    if not file_path.suffix.lower() == ".pdf":
        return "Error: file must be a .pdf"

    try:
        reader = PdfReader(str(file_path))
        texts = []
        for page in reader.pages:
            text = page.extract_text() or ""
            texts.append(text)
        combined = "\n".join(texts).strip()
        if not combined:
            return "PDF contains no extractable text (may be a scanned image)."
        if len(combined) > 4000:
            return combined[:4000] + "\n[truncated]"
        return combined
    except Exception as e:
        return f"PDF extraction failed: {e}"
